from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


TITLE_COLOR = RGBColor(17, 34, 68)
ACCENT = RGBColor(15, 122, 140)
ACCENT_DARK = RGBColor(9, 78, 90)
HIGHLIGHT = RGBColor(242, 165, 76)
TEXT = RGBColor(38, 47, 63)
MUTED = RGBColor(96, 110, 133)
LIGHT_BG = RGBColor(245, 247, 250)
WHITE = RGBColor(255, 255, 255)
SUCCESS = RGBColor(74, 136, 74)


def add_background(slide, color=LIGHT_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_top_band(slide, color=ACCENT_DARK):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.33), Inches(0.4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_title(slide, title, subtitle=None):
    add_background(slide)
    add_top_band(slide)
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.75), Inches(8.9), Inches(0.8))
    p = title_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Aptos Display"
    run.font.bold = True
    run.font.size = Pt(28)
    run.font.color.rgb = TITLE_COLOR
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.6), Inches(1.45), Inches(10.8), Inches(0.5))
        p2 = sub.text_frame.paragraphs[0]
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.name = "Aptos"
        run2.font.size = Pt(13)
        run2.font.color.rgb = MUTED


def add_bullets(slide, items, left=0.8, top=1.9, width=5.8, height=4.7, level0_size=20):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = text
        p.level = level
        p.font.name = "Aptos"
        p.font.size = Pt(level0_size - level * 2)
        p.font.color.rgb = TEXT
        p.space_after = Pt(8)
        if level == 0:
            p.bullet = True
    return box


def add_note_card(slide, title, body, left, top, width, height, fill_color=WHITE):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = RGBColor(220, 225, 232)
    tf = shape.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = title
    p1.font.name = "Aptos Display"
    p1.font.bold = True
    p1.font.size = Pt(18)
    p1.font.color.rgb = ACCENT_DARK
    p2 = tf.add_paragraph()
    p2.text = body
    p2.font.name = "Aptos"
    p2.font.size = Pt(13)
    p2.font.color.rgb = TEXT
    p2.space_before = Pt(8)
    return shape


def add_workflow_box(slide, left, top, width, height, title, body, fill):
    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = fill
    tf = box.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = title
    p1.font.name = "Aptos Display"
    p1.font.bold = True
    p1.font.size = Pt(18)
    p1.font.color.rgb = WHITE
    p1.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = body
    p2.font.name = "Aptos"
    p2.font.size = Pt(11)
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER
    return box


def add_arrow(slide, left, top, width):
    arrow = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.CHEVRON,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(0.42),
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = HIGHLIGHT
    arrow.line.color.rgb = HIGHLIGHT


def add_code_block(slide, title, code, left, top, width, height):
    label = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.35))
    p = label.text_frame.paragraphs[0]
    p.text = title
    p.font.name = "Aptos Display"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = ACCENT_DARK

    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top + 0.38),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(31, 36, 46)
    shape.line.color.rgb = RGBColor(31, 36, 46)
    tf = shape.text_frame
    tf.clear()
    p2 = tf.paragraphs[0]
    p2.text = code
    p2.font.name = "Consolas"
    p2.font.size = Pt(12)
    p2.font.color.rgb = RGBColor(227, 233, 239)
    p2.alignment = PP_ALIGN.LEFT
    return shape


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


# Slide 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, RGBColor(232, 244, 242))
band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
band.fill.solid()
band.fill.fore_color.rgb = RGBColor(232, 244, 242)
band.line.fill.background()
circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(8.9), Inches(0.9), Inches(3.2), Inches(3.2))
circle.fill.solid()
circle.fill.fore_color.rgb = ACCENT
circle.line.fill.background()
inner = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(9.55), Inches(1.55), Inches(1.9), Inches(1.9))
inner.fill.solid()
inner.fill.fore_color.rgb = HIGHLIGHT
inner.line.fill.background()
title = slide.shapes.add_textbox(Inches(0.75), Inches(1.15), Inches(6.8), Inches(1.5))
t1 = title.text_frame.paragraphs[0]
t1.text = "AI Assisted Coding"
t1.font.name = "Aptos Display"
t1.font.bold = True
t1.font.size = Pt(30)
t1.font.color.rgb = TITLE_COLOR
t2 = title.text_frame.add_paragraph()
t2.text = "Workflow, prompt patterns, and classroom-ready practices"
t2.font.name = "Aptos"
t2.font.size = Pt(18)
t2.font.color.rgb = MUTED
chip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.3), Inches(3.2), Inches(0.6))
chip.fill.solid()
chip.fill.fore_color.rgb = WHITE
chip.line.color.rgb = RGBColor(220, 225, 232)
ct = chip.text_frame.paragraphs[0]
ct.text = "Course deck | instructor-ready"
ct.font.name = "Aptos"
ct.font.size = Pt(16)
ct.font.color.rgb = ACCENT_DARK
ct.alignment = PP_ALIGN.CENTER


# Slide 2
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "Learning Goals", "What learners should be able to do by the end of the session")
add_bullets(
    slide,
    [
        "Explain how AI changes the software development workflow.",
        "Write prompts that produce better plans, code, tests, and refactors.",
        "Review AI output critically instead of accepting it blindly.",
        "Use a repeatable loop: ask, inspect, run, verify, refine.",
        "Apply guardrails for security, quality, and maintainability.",
    ],
)
add_note_card(slide, "Teaching Tip", "Position AI as a pair programmer, not an autopilot. Students should stay responsible for decisions and verification.", 7.15, 2.0, 5.1, 2.1, RGBColor(255, 249, 239))
add_note_card(slide, "Best Outcome", "Students leave with prompt templates they can reuse in labs, internships, and project work.", 7.15, 4.35, 5.1, 1.7)


# Slide 3
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "Where AI Fits In Coding", "AI helps across the lifecycle, but human judgment anchors the work")
cards = [
    ("Understand", "Clarify requirements, constraints, and edge cases.", ACCENT_DARK, 0.8),
    ("Design", "Compare approaches, data flow, APIs, and tradeoffs.", ACCENT, 3.25),
    ("Build", "Draft code, tests, migrations, and docs faster.", HIGHLIGHT, 5.7),
    ("Review", "Spot bugs, simplify logic, and strengthen tests.", SUCCESS, 8.15),
    ("Learn", "Explain unfamiliar code and generate examples.", RGBColor(113, 90, 163), 10.6),
]
for title_text, body, color, left in cards:
    add_workflow_box(slide, left, 2.45, 2.0, 2.0, title_text, body, color)

add_note_card(slide, "Key Message", "The strongest results come from short, iterative prompts with context, constraints, and a verification step.", 1.2, 5.2, 10.8, 1.3)


# Slide 4
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "Core Workflow", "A simple repeatable loop for AI-assisted coding")
workflow = [
    ("1. Define", "State the task,\ninputs, outputs,\nand constraints.", ACCENT_DARK, 0.55),
    ("2. Ask", "Request a plan,\npatch, test, or\nexplanation.", ACCENT, 3.0),
    ("3. Inspect", "Read the output,\ncheck assumptions,\ncompare options.", HIGHLIGHT, 5.45),
    ("4. Execute", "Run code, tests,\nlinters, and sample\ncases.", SUCCESS, 7.9),
    ("5. Refine", "Fix issues and\niterate with sharper\nprompts.", RGBColor(113, 90, 163), 10.35),
]
for index, (title_text, body, color, left) in enumerate(workflow):
    add_workflow_box(slide, left, 2.35, 2.05, 2.2, title_text, body, color)
    if index < len(workflow) - 1:
        add_arrow(slide, left + 2.1, 3.15, 0.55)
add_note_card(slide, "Instructor Script", "Emphasize that prompting is not a single message. It is a feedback loop that gets sharper as evidence comes in.", 1.1, 5.25, 11.0, 1.25)


# Slide 5
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "Prompt Formula", "A reusable structure that improves output quality")
add_note_card(slide, "Formula", "Role + task + context + constraints + expected output + verification request", 0.8, 1.75, 11.7, 0.95, RGBColor(235, 247, 248))
add_bullets(
    slide,
    [
        "Role: \"Act as a senior backend engineer\"",
        "Task: \"Add image upload validation\"",
        "Context: codebase, stack, existing files, user goal",
        "Constraints: no new dependencies, keep API unchanged, use async/await",
        "Expected output: patch, explanation, tests, commit message",
        "Verification: ask for edge cases, risks, and test scenarios",
    ],
    top=2.95,
    width=6.2,
)
add_note_card(slide, "Rule of Thumb", "If the answer is vague, the prompt was probably vague. Add specifics before asking again.", 7.45, 3.15, 4.8, 2.1, RGBColor(255, 249, 239))


# Slide 6
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "Prompt Examples", "Use task-specific prompts instead of one generic request")
add_code_block(
    slide,
    "Planning prompt",
    "Analyze this feature request and propose 2 implementation options.\nInclude tradeoffs, files likely to change, and a recommended path.",
    0.75,
    1.75,
    5.9,
    1.4,
)
add_code_block(
    slide,
    "Coding prompt",
    "Update the Express route in server.js to validate image type and size.\nDo not change the response schema. Show the exact patch.",
    6.8,
    1.75,
    5.8,
    1.4,
)
add_code_block(
    slide,
    "Testing prompt",
    "Write test cases for success, unsupported file type, oversized upload,\nand malformed request body. Explain what each test proves.",
    0.75,
    3.95,
    5.9,
    1.45,
)
add_code_block(
    slide,
    "Review prompt",
    "Review this code for bugs, security risks, and missing tests.\nPrioritize findings by severity and reference the relevant lines.",
    6.8,
    3.95,
    5.8,
    1.45,
)


# Slide 7
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "From Weak Prompt to Strong Prompt", "Teach students how to add the missing context")
add_note_card(slide, "Weak Prompt", "Create a login page.", 0.8, 2.0, 5.8, 1.25, RGBColor(255, 240, 240))
add_note_card(
    slide,
    "Why It Fails",
    "It does not say which framework, what fields are needed, how errors work, what style to follow, or what output format is expected.",
    0.8,
    3.55,
    5.8,
    1.85,
)
add_note_card(
    slide,
    "Strong Prompt",
    "Build a responsive React login page using our existing design system. Include email and password fields, inline validation, a loading state, and an error banner. Return the JSX and CSS changes only. Keep the current API contract unchanged.",
    6.8,
    2.0,
    5.5,
    3.4,
    RGBColor(237, 248, 239),
)


# Slide 8
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "Verification Checklist", "What students should do before trusting AI output")
add_bullets(
    slide,
    [
        "Does the code actually run in the target environment?",
        "Are imports, dependencies, and file paths correct?",
        "Did the model invent an API, config, or library behavior?",
        "Are edge cases, validation, and error handling covered?",
        "Do tests prove the expected behavior, not just happy paths?",
        "Can the code be simplified before merging?",
    ],
    width=6.0,
)
add_note_card(slide, "Mindset", "AI can accelerate delivery, but verification protects quality. Fast output is not the same as correct output.", 7.15, 2.1, 5.0, 1.9, RGBColor(255, 249, 239))
add_note_card(slide, "Good Practice", "Ask the model to list assumptions and likely risks. That often reveals hidden problems early.", 7.15, 4.4, 5.0, 1.55)


# Slide 9
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "Responsible Use", "Important guardrails for classroom and workplace adoption")
add_bullets(
    slide,
    [
        "Do not paste secrets, private keys, or confidential source code into public tools.",
        "Review licensing and usage policies before copying generated code into production.",
        "Treat security-sensitive code as high scrutiny: auth, payments, data access, and encryption.",
        "Keep a human in the loop for architecture decisions and final approval.",
        "Document what AI helped generate when your process requires traceability.",
    ],
    width=6.4,
)
add_note_card(slide, "Classroom Framing", "Use AI to learn faster and think better, not to skip understanding. Students should be able to explain any code they submit.", 7.2, 2.25, 4.9, 2.05, RGBColor(235, 247, 248))


# Slide 10
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "Hands-On Activity", "A lab structure for the course")
add_note_card(slide, "Task", "Students choose a small feature, bug fix, or refactor from a sample codebase.", 0.8, 1.9, 3.7, 1.3)
add_note_card(slide, "Prompt", "They write a planning prompt first, then a coding prompt with constraints.", 4.8, 1.9, 3.7, 1.3)
add_note_card(slide, "Verify", "They run the code, test assumptions, and document one improvement after feedback.", 8.8, 1.9, 3.7, 1.3)
add_bullets(
    slide,
    [
        "Deliverables: prompts used, generated output, issues found, final corrected version.",
        "Discussion: What did AI do well? Where did it need correction?",
        "Reflection: Which prompt detail most improved the result?",
    ],
    left=1.0,
    top=4.0,
    width=11.0,
    height=2.2,
    level0_size=18,
)


# Slide 11
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "Prompt Library for Students", "Reusable prompts they can adapt across assignments")
add_code_block(
    slide,
    "Explain code",
    "Explain what this file does, identify the main control flow,\nand describe the 3 most important functions in plain language.",
    0.7,
    1.65,
    4.0,
    1.55,
)
add_code_block(
    slide,
    "Refactor safely",
    "Refactor this function for readability without changing behavior.\nList any assumptions and suggest tests to confirm no regression.",
    4.7,
    1.65,
    4.0,
    1.55,
)
add_code_block(
    slide,
    "Debugging",
    "Given this error message and stack trace, suggest the most likely causes,\na step-by-step debugging path, and the minimal fix to try first.",
    8.7,
    1.65,
    3.9,
    1.55,
)
add_code_block(
    slide,
    "Documentation",
    "Write concise README steps for setup, run, and test.\nMention prerequisites and common troubleshooting notes.",
    0.7,
    4.25,
    5.8,
    1.35,
)
add_code_block(
    slide,
    "Code review",
    "Review this pull request for correctness, edge cases, and maintainability.\nReturn only findings ordered by severity.",
    6.7,
    4.25,
    5.8,
    1.35,
)


# Slide 12
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, RGBColor(240, 248, 247))
add_top_band(slide)
tb = slide.shapes.add_textbox(Inches(0.9), Inches(1.1), Inches(8.0), Inches(1.0))
p = tb.text_frame.paragraphs[0]
p.text = "Closing Takeaways"
p.font.name = "Aptos Display"
p.font.bold = True
p.font.size = Pt(28)
p.font.color.rgb = TITLE_COLOR
add_bullets(
    slide,
    [
        "AI-assisted coding works best as an iterative collaboration.",
        "Prompt quality improves with context, constraints, and clear expected output.",
        "Verification is essential: run, test, review, and refine.",
        "Students gain the most when they can explain and defend the final solution.",
    ],
    left=1.0,
    top=2.3,
    width=7.2,
    height=3.0,
)
quote = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.6), Inches(1.9), Inches(3.8), Inches(2.8))
quote.fill.solid()
quote.fill.fore_color.rgb = ACCENT
quote.line.fill.background()
tf = quote.text_frame
q1 = tf.paragraphs[0]
q1.text = "Best practice:"
q1.font.name = "Aptos Display"
q1.font.bold = True
q1.font.size = Pt(18)
q1.font.color.rgb = WHITE
q2 = tf.add_paragraph()
q2.text = "Ask better.\nCheck harder.\nShip smarter."
q2.font.name = "Aptos Display"
q2.font.bold = True
q2.font.size = Pt(24)
q2.font.color.rgb = WHITE
q2.alignment = PP_ALIGN.CENTER


output = "AI_Assisted_Coding_Workflow_and_Prompts.pptx"
prs.save(output)
print(output)
